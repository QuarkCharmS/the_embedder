"""
Minimalistic file content reader.
Returns file content as string or empty string if unsupported/error.
"""

import pathlib
import csv
from .chunker_config import EXT_TO_LANG

# Document format imports - fail silently
try:
    import chardet
    HAS_CHARDET = True
except ImportError:
    HAS_CHARDET = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from odf import text, teletype
    from odf.opendocument import load
    from odf.table import TableRow, TableCell
    HAS_ODF = True
except ImportError:
    HAS_ODF = False

try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from striprtf.striprtf import rtf_to_text
    HAS_RTF = True
except ImportError:
    HAS_RTF = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

try:
    import markdown
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False


def _detect_encoding(file_path):
    """Detect file encoding."""
    if HAS_CHARDET:
        try:
            with open(file_path, 'rb') as f:
                result = chardet.detect(f.read(10000))
                return result['encoding'] or 'utf-8'
        except Exception as e:
            import sys
            print(f"[WARN] chardet failed for {file_path}: {e}", file=sys.stderr)

    for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                f.read(1000)
            return encoding
        except Exception:
            continue
    return 'utf-8'


def _read_text_file(file_path):
    """Read text file with encoding detection."""
    encoding = _detect_encoding(file_path)
    try:
        with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
            return f.read()
    except Exception as e:
        import sys
        print(f"[ERROR] _read_text_file failed for {file_path}: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return ""


def read_file_content(file_path):
    """
    Read file content and return as string.
    Returns empty string if file is unsupported, symlink, or error occurs.
    """
    try:
        file_path = pathlib.Path(file_path)
        
        # Check if file exists and is not a symlink
        if not file_path.exists() or file_path.is_symlink():
            return ""
        
        ext = file_path.suffix.lower()
        filename = file_path.name
        
        # Document formats first
        if ext == '.pdf' and HAS_PDF:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                return '\n'.join(page.extract_text() for page in reader.pages)
        
        elif ext in ['.docx', '.doc'] and HAS_DOCX:
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        
        elif ext == '.odt' and HAS_ODF:
            doc = load(file_path)
            allparas = doc.getElementsByType(text.P)
            return '\n'.join(teletype.extractText(para) for para in allparas if teletype.extractText(para).strip())
        
        elif ext == '.rtf' and HAS_RTF:
            return rtf_to_text(_read_text_file(file_path))
        
        elif ext in ['.xlsx', '.xls'] and HAS_EXCEL:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            content = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    if any(cell for cell in row if cell is not None):
                        content.append(' | '.join(str(cell) if cell else '' for cell in row))
            return '\n'.join(content)
        
        elif ext == '.ods' and HAS_ODF:
            doc = load(file_path)
            rows = doc.spreadsheet.getElementsByType(TableRow)
            content = []
            for row in rows:
                cells = row.getElementsByType(TableCell)
                row_data = []
                for cell in cells:
                    paragraphs = cell.getElementsByType(text.P)
                    cell_text = "".join(teletype.extractText(p) for p in paragraphs)
                    row_data.append(cell_text.strip())
                if any(cell.strip() for cell in row_data):
                    content.append(' | '.join(row_data))
            return '\n'.join(content)
        
        elif ext in ['.pptx', '.ppt'] and HAS_PPTX:
            prs = Presentation(file_path)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text.strip())
            return '\n'.join(content)
        
        elif ext == '.csv':
            encoding = _detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                return '\n'.join(' | '.join(row) for row in csv.reader(f))
        
        elif ext in ['.html', '.htm'] and HAS_BS4:
            content = _read_text_file(file_path)
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        
        elif ext in ['.md', '.markdown'] and HAS_MARKDOWN and HAS_BS4:
            content = _read_text_file(file_path)
            html = markdown.markdown(content)
            return BeautifulSoup(html, 'html.parser').get_text()
        
        elif ext == '.xml' and HAS_BS4:
            content = _read_text_file(file_path)
            return BeautifulSoup(content, 'xml').get_text()
        
        # Check if it's a known code/text file from chunker_config
        elif ext in EXT_TO_LANG or filename in EXT_TO_LANG:
            return _read_text_file(file_path)
        
        # Try as text file for common extensions
        elif ext in ['.txt', '.text', '.log', '.ini', '.cfg', '.conf', '.env', '.properties']:
            return _read_text_file(file_path)
        
        # Default: empty string for unsupported
        return ""

    except Exception as e:
        # Log the error instead of silently returning empty
        import sys
        print(f"[ERROR] Failed to read file {file_path}: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return ""
